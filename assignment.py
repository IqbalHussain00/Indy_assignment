import os
from pptx import Presentation
from wand import image
from wand.image import Image
from wand.display import display
from wand.drawing import Drawing
from pptx.util import *

# this function watermark the image with logo
def watermarking(img,logo):

    image=os.getcwd()+'/files/'+img

    with Image(filename=logo) as logo_img:
        logo_img.resize(1000,500)
        with Image(filename=image) as image:        
            image.composite_channel('all_channels',logo_img,'dissolve',0,0)
            image.save(filename=os.getcwd()+'/files/'+img)



def convert_into_ppt(pr1,path_imgs,logo):

    logo=os.getcwd()+'/files/'+logo
    i=0                 #this variable is used to title and subtitle number on page of ppt 
    for img in path_imgs:
        watermarking(img,logo)
        i+=1
        layout = pr1.slide_layouts[1] 
        slide = pr1.slides.add_slide(layout)
        title = slide.shapes.title.text = 'Sample Title '+ str(i)
        sub = slide.placeholders[1].text = 'Sample Subtitle '+ str(i)

        im = Image(filename=os.getcwd()+'/files/'+img)
        width, height = im.size

        pic = slide.shapes.add_picture(os.getcwd()+'/files/'+img,Inches(1), Inches(2.5),width=Inches(width/1500), height=Inches(height/1500))
        

    pr1.save("files/Indycium.pptx")
        




path_imgs=os.listdir(os.getcwd()+'/files')[0:5]
logo=os.listdir(os.getcwd()+'/files')[5]


pr1=Presentation()                      # initilize presentation object

convert_into_ppt(pr1,path_imgs,logo)

os.startfile("files\Indycium.pptx")  
